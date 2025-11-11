"""
PowerPoint export functionality for narratives and dashboards.
"""

from typing import Dict, Any, Optional
from io import BytesIO
import logging

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. Install with: pip install python-pptx")


def create_powerpoint_from_narrative(
    narrative: Dict[str, Any],
    query: str,
    search_metadata: Optional[Dict[str, Any]] = None
) -> BytesIO:
    """
    Create a PowerPoint presentation from a narrative.
    
    Args:
        narrative: Narrative data with markdown content
        query: Original search query
        search_metadata: Optional search metadata
        
    Returns:
        BytesIO object containing the PowerPoint file
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"Search Analysis: {query}"
    subtitle.text = "Generated Narrative Report"
    
    # Get narrative content
    narrative_text = narrative.get("markdown", "")
    sections = narrative.get("sections", {})
    
    # If we have sections, create slides for each
    if sections:
        for section_name, section_content in sections.items():
            # Create a new slide for each section
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            # Format section name
            section_title = section_name.replace("_", " ").title()
            title_shape.text = section_title
            
            # Add content as bullet points
            tf = body_shape.text_frame
            tf.text = section_content[:200]  # First 200 chars
            
            # Split into paragraphs if content is long
            if len(section_content) > 200:
                paragraphs = section_content.split("\n")
                for para in paragraphs[1:6]:  # Limit to 5 more paragraphs
                    if para.strip():
                        p = tf.add_paragraph()
                        p.text = para.strip()[:100]
                        p.level = 0
    else:
        # Create slides from markdown content
        lines = narrative_text.split("\n")
        current_slide = None
        current_title = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check for headers
            if line.startswith("#"):
                # Save previous slide if exists
                if current_slide:
                    pass  # Previous slide is already added
                
                # Create new slide
                bullet_slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = current_slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                # Extract title
                current_title = line.lstrip("#").strip()
                title_shape.text = current_title
                tf = body_shape.text_frame
                tf.text = ""  # Clear default text
            elif current_slide:
                # Add content to current slide
                shapes = current_slide.shapes
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                
                if line.startswith("- ") or line.startswith("* "):
                    # Bullet point
                    p = tf.add_paragraph()
                    p.text = line[2:].strip()
                    p.level = 0
                else:
                    # Regular paragraph
                    if not tf.text:
                        tf.text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                        p.level = 0
    
    # Add metadata slide if available
    if search_metadata:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Search Metadata"
        tf = body_shape.text_frame
        
        tf.text = f"Query: {query}"
        p = tf.add_paragraph()
        p.text = f"Sources Queried: {search_metadata.get('sources_queried', 0)}"
        p.level = 0
        p = tf.add_paragraph()
        p.text = f"Sources Successful: {search_metadata.get('sources_successful', 0)}"
        p.level = 0
        p = tf.add_paragraph()
        p.text = f"Execution Time: {search_metadata.get('execution_time_ms', 0):.2f} ms"
        p.level = 0
    
    # Save to BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output


def create_powerpoint_from_dashboard(
    dashboard: Dict[str, Any],
    query: str,
    search_metadata: Optional[Dict[str, Any]] = None
) -> BytesIO:
    """
    Create a PowerPoint presentation from a dashboard specification.
    
    Args:
        dashboard: Dashboard specification
        query: Original search query
        search_metadata: Optional search metadata
        
    Returns:
        BytesIO object containing the PowerPoint file
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    dashboard_title = dashboard.get("specification", {}).get("title", f"Dashboard: {query}")
    title.text = dashboard_title
    subtitle.text = dashboard.get("specification", {}).get("description", "Generated Dashboard Report")
    
    spec = dashboard.get("specification", {})
    
    # Metrics slide
    metrics = spec.get("metrics", [])
    if metrics:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Key Metrics"
        tf = body_shape.text_frame
        
        for metric in metrics[:10]:  # Limit to 10 metrics
            label = metric.get("label", "")
            value = metric.get("value", "")
            format_type = metric.get("format", "number")
            
            if format_type == "percentage" and isinstance(value, (int, float)):
                value_str = f"{value * 100:.1f}%"
            elif format_type == "currency" and isinstance(value, (int, float)):
                value_str = f"${value:,.2f}"
            else:
                value_str = str(value)
            
            p = tf.add_paragraph()
            p.text = f"{label}: {value_str}"
            p.level = 0
    
    # Insights slide
    insights = spec.get("insights", [])
    if insights:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Key Insights"
        tf = body_shape.text_frame
        
        for insight in insights[:10]:  # Limit to 10 insights
            p = tf.add_paragraph()
            p.text = insight
            p.level = 0
    
    # Charts slide
    charts = spec.get("charts", [])
    if charts:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Recommended Visualizations"
        tf = body_shape.text_frame
        
        for chart in charts[:10]:  # Limit to 10 charts
            chart_title = chart.get("title", "Untitled Chart")
            chart_type = chart.get("type", "unknown")
            data_source = chart.get("data_source", "unknown")
            
            p = tf.add_paragraph()
            p.text = f"{chart_title} ({chart_type})"
            p.level = 0
            p = tf.add_paragraph()
            p.text = f"  Data Source: {data_source}"
            p.level = 1
    
    # Add metadata slide if available
    if search_metadata:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Search Metadata"
        tf = body_shape.text_frame
        
        tf.text = f"Query: {query}"
        p = tf.add_paragraph()
        p.text = f"Sources Queried: {search_metadata.get('sources_queried', 0)}"
        p.level = 0
        p = tf.add_paragraph()
        p.text = f"Sources Successful: {search_metadata.get('sources_successful', 0)}"
        p.level = 0
        p = tf.add_paragraph()
        p.text = f"Execution Time: {search_metadata.get('execution_time_ms', 0):.2f} ms"
        p.level = 0
    
    # Save to BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output


def create_powerpoint_from_narrative_and_dashboard(
    narrative: Dict[str, Any],
    dashboard: Dict[str, Any],
    query: str,
    search_metadata: Optional[Dict[str, Any]] = None
) -> BytesIO:
    """
    Create a combined PowerPoint presentation from narrative and dashboard.
    
    Args:
        narrative: Narrative data
        dashboard: Dashboard specification
        query: Original search query
        search_metadata: Optional search metadata
        
    Returns:
        BytesIO object containing the PowerPoint file
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    dashboard_title = dashboard.get("specification", {}).get("title", f"Report: {query}")
    title.text = dashboard_title or f"Search Analysis: {query}"
    subtitle.text = "Comprehensive Analysis Report"
    
    # Narrative section
    narrative_text = narrative.get("markdown", "")
    sections = narrative.get("sections", {})
    
    if sections:
        # Section divider
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Narrative Analysis"
        
        for section_name, section_content in sections.items():
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            section_title = section_name.replace("_", " ").title()
            title_shape.text = section_title
            
            tf = body_shape.text_frame
            paragraphs = section_content.split("\n")
            for para in paragraphs[:8]:  # Limit paragraphs
                if para.strip():
                    p = tf.add_paragraph()
                    p.text = para.strip()[:150]
                    p.level = 0
    
    # Dashboard section
    spec = dashboard.get("specification", {})
    
    # Section divider
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Dashboard & Metrics"
    
    # Metrics
    metrics = spec.get("metrics", [])
    if metrics:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Key Metrics"
        tf = body_shape.text_frame
        
        for metric in metrics[:10]:
            label = metric.get("label", "")
            value = metric.get("value", "")
            format_type = metric.get("format", "number")
            
            if format_type == "percentage" and isinstance(value, (int, float)):
                value_str = f"{value * 100:.1f}%"
            elif format_type == "currency" and isinstance(value, (int, float)):
                value_str = f"${value:,.2f}"
            else:
                value_str = str(value)
            
            p = tf.add_paragraph()
            p.text = f"{label}: {value_str}"
            p.level = 0
    
    # Insights
    insights = spec.get("insights", [])
    if insights:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Key Insights"
        tf = body_shape.text_frame
        
        for insight in insights[:10]:
            p = tf.add_paragraph()
            p.text = insight
            p.level = 0
    
    # Charts
    charts = spec.get("charts", [])
    if charts:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Recommended Visualizations"
        tf = body_shape.text_frame
        
        for chart in charts[:10]:
            chart_title = chart.get("title", "Untitled Chart")
            chart_type = chart.get("type", "unknown")
            
            p = tf.add_paragraph()
            p.text = f"{chart_title} ({chart_type})"
            p.level = 0
    
    # Metadata slide
    if search_metadata:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Search Metadata"
        tf = body_shape.text_frame
        
        tf.text = f"Query: {query}"
        p = tf.add_paragraph()
        p.text = f"Sources Queried: {search_metadata.get('sources_queried', 0)}"
        p.level = 0
        p = tf.add_paragraph()
        p.text = f"Sources Successful: {search_metadata.get('sources_successful', 0)}"
        p.level = 0
        p = tf.add_paragraph()
        p.text = f"Execution Time: {search_metadata.get('execution_time_ms', 0):.2f} ms"
        p.level = 0
    
    # Save to BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output

